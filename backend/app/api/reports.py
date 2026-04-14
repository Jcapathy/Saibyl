from __future__ import annotations

import asyncio
import io
import json
import re

import structlog
from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel

from app.core.auth import get_current_org
from app.core.database import get_supabase_admin
from app.services.intelligence.report_agent import clean_report_output, get_report_progress, strip_react_artifacts
from app.services.intelligence.report_chat import chat_with_report
from app.workers.report_tasks import run_generate_report

log = structlog.get_logger()


def _compute_polarization(events: list[dict]) -> dict:
    """Compute polarization metrics from simulation event sentiment values.

    Uses per-agent sentiment at the final round to compute the extreme-to-moderate
    ratio.  Returns controversy_score (0-1), polarization_ratio (str like "2.7:1"),
    and valence_switching_pct (int 0-100).
    """
    if not events:
        return {"controversy_score": None, "polarization_ratio": None, "valence_switching_pct": None}

    # Find the maximum round number (final round)
    max_round = 0
    for e in events:
        rn = e.get("round_number") or 0
        if rn > max_round:
            max_round = rn

    # Collect per-agent sentiment at the final round (deduplicated: last event wins)
    agent_sentiments: dict[str, float] = {}
    all_sentiments: list[float] = []
    for e in events:
        md = e.get("metadata") or {}
        s = md.get("sentiment")
        if s is None:
            continue
        try:
            val = float(s)
        except (ValueError, TypeError):
            continue
        all_sentiments.append(val)
        rn = e.get("round_number") or 0
        if rn == max_round and e.get("agent_id"):
            agent_sentiments[e["agent_id"]] = val

    # Use per-agent final-round sentiments for ratio; fall back to all sentiments
    sentiments = list(agent_sentiments.values()) if agent_sentiments else all_sentiments
    if not sentiments:
        return {"controversy_score": None, "polarization_ratio": None, "valence_switching_pct": None}

    # Extreme-to-moderate ratio: |sentiment| > 0.5 vs |sentiment| <= 0.5
    extreme = sum(1 for s in sentiments if abs(s) > 0.5)
    moderate = max(sum(1 for s in sentiments if abs(s) <= 0.5), 1)
    ratio = round(extreme / moderate, 1)

    # Valence switching: % of consecutive pairs that cross the zero line (all events)
    switches = 0
    for i in range(1, len(all_sentiments)):
        if (all_sentiments[i] > 0) != (all_sentiments[i - 1] > 0):
            switches += 1
    switching_pct = round(switches / max(len(all_sentiments) - 1, 1) * 100)

    # Normalize ratio to 0-1 scale (ratio of 5:1+ saturates at 1.0)
    controversy_score = round(min(1.0, ratio / 5.0), 2)

    return {
        "controversy_score": controversy_score,
        "polarization_ratio": f"{ratio}:1",
        "valence_switching_pct": switching_pct,
    }


async def _safe_task(coro, name: str):
    try:
        await coro
    except Exception:
        log.exception("background_task_failed", task=name)

router = APIRouter(tags=["reports"])


# ---------------------------------------------------------------------------
# Request bodies
# ---------------------------------------------------------------------------

class GenerateReportBody(BaseModel):
    simulation_id: str
    variant: str = "a"
    evidence_depth: str = "deep"  # shallow, standard, deep, exhaustive
    max_sections: int | None = None


class ExportReportBody(BaseModel):
    format: str  # "json", "pdf", "pptx"


class ChatBody(BaseModel):
    message: str
    history: list[dict[str, str]] | None = None


# ---------------------------------------------------------------------------
# Endpoints
# ---------------------------------------------------------------------------

@router.post("/generate")
async def generate_report(body: GenerateReportBody, auth: dict = Depends(get_current_org)):
    """Trigger report generation for a simulation."""
    log.info("generate_report", simulation_id=body.simulation_id, org_id=auth["org_id"])
    admin = get_supabase_admin()

    # Verify simulation belongs to org
    sim = (
        admin.table("simulations")
        .select("id")
        .eq("id", body.simulation_id)
        .eq("organization_id", auth["org_id"])
        .single()
        .execute()
    )
    if not sim.data:
        raise HTTPException(status_code=404, detail="Simulation not found")

    asyncio.create_task(_safe_task(
        run_generate_report(body.simulation_id, body.variant, body.evidence_depth, body.max_sections),
        "generate_report",
    ))
    return {"status": "started"}


@router.get("/by-simulation/{sim_id}")
async def get_reports_by_simulation(sim_id: str, auth: dict = Depends(get_current_org)):
    """Get the latest report for a simulation, with sections embedded."""
    log.info("get_reports_by_simulation", simulation_id=sim_id, org_id=auth["org_id"])
    admin = get_supabase_admin()

    # Verify simulation belongs to org
    sim = (
        admin.table("simulations")
        .select("id, project_id")
        .eq("id", sim_id)
        .eq("organization_id", auth["org_id"])
        .single()
        .execute()
    )
    if not sim.data:
        raise HTTPException(status_code=404, detail="Simulation not found")

    result = (
        admin.table("reports")
        .select("*")
        .eq("simulation_id", sim_id)
        .order("created_at", desc=True)
        .limit(1)
        .execute()
    )
    if not result.data:
        raise HTTPException(status_code=404, detail="No report found for this simulation")

    report = result.data[0]

    # Load sections from report_sections table
    sections_result = (
        admin.table("report_sections")
        .select("title, content")
        .eq("report_id", report["id"])
        .order("section_index")
        .execute()
    )

    # Compute polarization metrics from simulation events
    events = (
        admin.table("simulation_events")
        .select("metadata, round_number, agent_id")
        .eq("simulation_id", sim_id)
        .limit(2000)
        .execute()
    ).data or []
    polarization = _compute_polarization(events)

    # Fetch source documents for the simulation's project
    source_documents: list[dict] = []
    project_id = sim.data.get("project_id")
    if project_id:
        docs = (
            admin.table("documents")
            .select("id, filename, file_type, storage_path, file_size_bytes")
            .eq("project_id", project_id)
            .eq("processing_status", "complete")
            .order("created_at")
            .limit(5)
            .execute()
        ).data or []
        for doc in docs:
            try:
                file_bytes = admin.storage.from_("project-media").download(doc["storage_path"])
                text = file_bytes.decode("utf-8", errors="replace")
                word_count = len(text.split())
                # Truncate to first ~500 words if over 2000 chars
                if len(text) > 2000:
                    words = text.split()[:500]
                    text = " ".join(words)
                    text += f"\n\n[Full source material: {word_count:,} words total]"
                source_documents.append({
                    "filename": doc["filename"],
                    "file_type": doc["file_type"],
                    "word_count": word_count,
                    "text": text,
                })
            except Exception:
                log.warning("source_doc_fetch_failed", doc_id=doc["id"])
                source_documents.append({
                    "filename": doc["filename"],
                    "file_type": doc["file_type"],
                    "word_count": 0,
                    "text": "[Document could not be loaded]",
                })

    # Return shape the frontend expects — strip any ReACT artifacts from content
    return {
        "id": report["id"],
        "simulation_id": report["simulation_id"],
        "status": report.get("status"),
        "sections": [
            {"title": s["title"], "content": strip_react_artifacts(s.get("content") or "")}
            for s in (sections_result.data or [])
        ],
        "full_markdown": strip_react_artifacts(report.get("markdown_content") or ""),
        "polarization": polarization,
        "source_documents": source_documents,
    }


@router.get("/{id}")
async def get_report(id: str, auth: dict = Depends(get_current_org)):
    """Get full report."""
    log.info("get_report", report_id=id)
    admin = get_supabase_admin()
    result = (
        admin.table("reports")
        .select("*, simulations!inner(organization_id)")
        .eq("id", id)
        .eq("simulations.organization_id", auth["org_id"])
        .single()
        .execute()
    )
    if not result.data:
        raise HTTPException(status_code=404, detail="Report not found")
    return result.data


@router.get("/{id}/sections")
async def list_report_sections(id: str, auth: dict = Depends(get_current_org)):
    """List sections of a report."""
    log.info("list_report_sections", report_id=id)
    admin = get_supabase_admin()

    # Verify report belongs to org via simulation join
    report = (
        admin.table("reports")
        .select("id, simulations!inner(organization_id)")
        .eq("id", id)
        .eq("simulations.organization_id", auth["org_id"])
        .single()
        .execute()
    )
    if not report.data:
        raise HTTPException(status_code=404, detail="Report not found")

    result = (
        admin.table("report_sections")
        .select("*")
        .eq("report_id", id)
        .order("section_index")
        .execute()
    )
    return result.data


@router.get("/{id}/progress")
async def report_progress(id: str, auth: dict = Depends(get_current_org)):
    """Get report generation progress."""
    log.info("report_progress", report_id=id)
    admin = get_supabase_admin()

    # Verify report belongs to org via simulation join
    report = (
        admin.table("reports")
        .select("id, simulations!inner(organization_id)")
        .eq("id", id)
        .eq("simulations.organization_id", auth["org_id"])
        .single()
        .execute()
    )
    if not report.data:
        raise HTTPException(status_code=404, detail="Report not found")

    progress = get_report_progress(id)
    return progress.model_dump()


@router.post("/{id}/chat")
async def chat_with_report_endpoint(id: str, body: ChatBody, auth: dict = Depends(get_current_org)):
    """Chat with a report using tool-augmented answers."""
    log.info("chat_with_report", report_id=id)
    admin = get_supabase_admin()

    # Verify report belongs to org via simulation join
    report = (
        admin.table("reports")
        .select("id, simulations!inner(organization_id)")
        .eq("id", id)
        .eq("simulations.organization_id", auth["org_id"])
        .single()
        .execute()
    )
    if not report.data:
        raise HTTPException(status_code=404, detail="Report not found")

    result = await chat_with_report(id, body.message, body.history)
    return result.model_dump()


@router.post("/{id}/export")
async def export_report(id: str, body: ExportReportBody, auth: dict = Depends(get_current_org)):
    """Export report as JSON, PDF, or PPTX — returns file download directly."""
    log.info("export_report", report_id=id, format=body.format)
    admin = get_supabase_admin()

    report = (
        admin.table("reports")
        .select("*, simulations!inner(name, organization_id)")
        .eq("id", id)
        .eq("simulations.organization_id", auth["org_id"])
        .single()
        .execute()
    )
    if not report.data:
        raise HTTPException(status_code=404, detail="Report not found")

    rpt = report.data
    markdown = strip_react_artifacts(rpt.get("markdown_content") or "")
    title = rpt.get("title") or "Report"
    safe_title = re.sub(r"[^a-zA-Z0-9_\- ]", "", title)[:80].strip()

    # Load sections
    sections = (
        admin.table("report_sections")
        .select("title, content, section_index")
        .eq("report_id", id)
        .order("section_index")
        .execute()
    ).data or []

    if body.format == "json":
        payload = {
            "id": rpt["id"],
            "title": title,
            "simulation_id": rpt["simulation_id"],
            "status": rpt.get("status"),
            "sections": [{"title": s["title"], "content": strip_react_artifacts(s.get("content") or "")} for s in sections],
            "full_markdown": markdown,
        }
        buf = io.BytesIO(json.dumps(payload, indent=2, ensure_ascii=False).encode("utf-8"))
        return StreamingResponse(
            buf,
            media_type="application/json",
            headers={"Content-Disposition": f'attachment; filename="{safe_title}.json"'},
        )

    if body.format == "pdf":
        import markdown as md
        from weasyprint import HTML

        html_body = md.markdown(markdown, extensions=["tables", "fenced_code"])
        full_html = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8">
<style>
  body {{ font-family: 'Helvetica Neue', Arial, sans-serif; max-width: 800px; margin: 40px auto; padding: 0 20px; color: #1a1a1a; line-height: 1.6; }}
  h1 {{ color: #2d2d7f; border-bottom: 2px solid #5B5FEE; padding-bottom: 8px; }}
  h2 {{ color: #3d3d9f; margin-top: 32px; }}
  blockquote {{ border-left: 3px solid #5B5FEE; padding-left: 16px; color: #555; margin: 16px 0; }}
  table {{ border-collapse: collapse; width: 100%; margin: 16px 0; }}
  th, td {{ border: 1px solid #ddd; padding: 8px 12px; text-align: left; }}
  th {{ background: #f5f5ff; }}
  code {{ background: #f4f4f4; padding: 2px 6px; border-radius: 3px; font-size: 0.9em; }}
</style>
</head><body>{html_body}</body></html>"""
        pdf_bytes = HTML(string=full_html).write_pdf()
        buf = io.BytesIO(pdf_bytes)
        return StreamingResponse(
            buf,
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{safe_title}.pdf"'},
        )

    if body.format == "pptx":
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = "Generated by Saibyl"

        # Section slides
        for sec in sections:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = sec["title"]
            content = strip_react_artifacts(sec.get("content") or "")
            # Strip markdown formatting for plain text slides
            plain = re.sub(r"[#*_`>]", "", content)
            tf = slide.placeholders[1].text_frame
            tf.text = plain[:3000]
            for para in tf.paragraphs:
                para.font.size = Pt(14)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return StreamingResponse(
            buf,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{safe_title}.pptx"'},
        )

    raise HTTPException(status_code=400, detail=f"Unsupported format: {body.format}. Use json, pdf, or pptx.")


@router.delete("/{id}")
async def delete_report(id: str, auth: dict = Depends(get_current_org)):
    """Delete a report and its sections."""
    log.info("delete_report", report_id=id)
    admin = get_supabase_admin()

    # Verify report belongs to org via simulation join
    report = (
        admin.table("reports")
        .select("id, simulations!inner(organization_id)")
        .eq("id", id)
        .eq("simulations.organization_id", auth["org_id"])
        .single()
        .execute()
    )
    if not report.data:
        raise HTTPException(status_code=404, detail="Report not found")

    # Delete sections first, then report
    admin.table("report_sections").delete().eq("report_id", id).execute()
    admin.table("reports").delete().eq("id", id).execute()

    return {"detail": "Report deleted"}
