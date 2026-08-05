# PUBLIC INTERFACE
# ─────────────────────────────────────────────────────────
# export_report_pptx(report_id) -> bytes
# ─────────────────────────────────────────────────────────
"""The deck version of the report, built on the same artifact as the PDF.

This shipped with zero charts. All three of its `simulation_analytics` calls
read keys the refactored tool no longer returns — `persona_events`,
`sentiment_curve`-as-a-dict, `platform_events` — and the whole chart block sat
inside a `try` that logged a warning and moved on, so every deck came out with
the slides missing and nothing said so.

It also read its headline figure by running a regular expression over the report
markdown looking for "Sentiment Trajectory". That is the scraping the
measurement layer exists to replace: it reads whatever the model happened to
write, including a number the model invented.

Both are gone. The deck now reads the same validated artifact the PDF does, and
carries the same three rules: unmeasured is absent rather than zero, no mean is
shown without its band, and an unresolved comparison is stated as unresolved
rather than presented as a ranking.
"""
from __future__ import annotations

import io
from datetime import UTC, datetime
from uuid import UUID

import structlog
from pptx import Presentation
from pptx.util import Inches, Pt

from app.core.database import get_supabase_admin
from app.services.export.chart_renderer import (
    ChartRow,
    render_interval_rows_png,
    render_sentiment_arc_png,
)
from app.services.export.report_document import (
    COHORT_NAMES,
    load_artifact,
    platform_label,
)
from app.services.intelligence.analysis_schema import SimulationAnalysis
from app.services.intelligence.report_agent import clean_report_output

logger = structlog.get_logger()

BODY_SIZE = Pt(14)
NOTE_SIZE = Pt(11)


def _title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _bullet_slide(
    prs: Presentation, title: str, bullets: list[str], *, note: str = ""
) -> None:
    bullets = [b for b in bullets if b and b.strip()]
    if not bullets:
        return
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    frame = slide.placeholders[1].text_frame
    frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.font.size = BODY_SIZE
    if note:
        paragraph = frame.add_paragraph()
        paragraph.text = note
        paragraph.font.size = NOTE_SIZE
        paragraph.font.italic = True


def _chart_slide(prs: Presentation, title: str, png: bytes) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.size = Pt(24)
    paragraph.font.bold = True
    slide.shapes.add_picture(io.BytesIO(png), Inches(0.8), Inches(1.2), width=Inches(11.7))


def _interval_rows(slices, label_of) -> list[ChartRow]:
    """Measured slices only. `n == 0` never reaches a chart."""
    return [
        ChartRow(
            label=label_of(slice_),
            mean=slice_.valence.mean,
            lower=slice_.valence.lower,
            upper=slice_.valence.upper,
            n=slice_.valence.n,
        )
        for slice_ in slices
        if slice_.valence.n > 0
    ]


def _headline_bullets(analysis: SimulationAnalysis) -> list[str]:
    headline = analysis.headline
    quality = analysis.quality
    bullets: list[str] = []
    if headline.valence.n > 0:
        bullets.append(
            f"Overall valence {headline.valence.mean:+.2f} "
            f"(95% CI {headline.valence.lower:+.2f} to "
            f"{headline.valence.upper:+.2f}, {headline.valence.n} agents)"
        )
    bullets.append(
        f"Stance of measured events: {headline.stance.support_pct:.0f}% support, "
        f"{headline.stance.oppose_pct:.0f}% oppose, "
        f"{headline.stance.undecided_pct:.0f}% undecided"
    )
    bullets.append(
        f"Trajectory {headline.trajectory} "
        f"({headline.trajectory_delta:+.2f} first round to last; reported flat "
        f"unless the intervals separate)"
    )
    bullets.append(
        f"Coverage {quality.coverage_pct:.1f}% — {quality.events_measured:,} of "
        f"{quality.events_total:,} events scored across {quality.agents_active} "
        f"active agents; confidence {quality.confidence}"
    )
    return bullets


async def export_report_pptx(report_id: str | UUID) -> bytes:
    """Generate the deck for one report."""
    admin = get_supabase_admin()
    report_id = str(report_id)

    report = (
        admin.table("reports").select("*").eq("id", report_id).limit(1).execute()
    ).data[0]
    simulation_id = str(report["simulation_id"])
    simulation = (
        admin.table("simulations").select("*").eq("id", simulation_id).limit(1).execute()
    ).data[0]
    sections = (
        admin.table("report_sections")
        .select("section_index, title, content")
        .eq("report_id", report_id)
        .order("section_index")
        .execute()
    ).data or []

    from app.services.intelligence.analysis_builder import get_analysis

    stored = get_analysis(simulation_id) or {}
    artifact = stored.get("artifact") if stored.get("build_status") == "complete" else None
    analysis, absence_reason = load_artifact(
        stored.get("schema_version") if artifact else None, artifact
    )

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    _title_slide(
        prs,
        report.get("title") or "Predictive intelligence report",
        f"{simulation.get('name') or 'Simulation'}\n"
        f"{datetime.now(UTC).strftime('%d %B %Y')}",
    )

    if analysis is None:
        _bullet_slide(prs, "No measured figures", [absence_reason])
    else:
        _bullet_slide(
            prs, "What the swarm did", _headline_bullets(analysis),
            note="Confidence intervals are computed across agents, not events.",
        )

        if analysis.quality.caveats:
            _bullet_slide(
                prs, "What this run can and cannot show",
                list(analysis.quality.caveats),
            )

        scoreboard = analysis.scoreboard
        if scoreboard is not None:
            if scoreboard.winner_variant_key:
                winner = next(
                    (
                        v for v in scoreboard.variants
                        if v.variant_key == scoreboard.winner_variant_key
                    ),
                    None,
                )
                name = (winner.label or winner.variant_key) if winner else "—"
                _bullet_slide(
                    prs, "Verdict",
                    [f"Winner: {name}", scoreboard.verdict],
                )
            else:
                _bullet_slide(
                    prs, "Verdict: no winner",
                    [
                        scoreboard.verdict
                        or "No variant separated from the others at 95% confidence.",
                        "The arenas below are in display order. That order is "
                        "not a ranking — where the intervals overlap this run "
                        "does not establish that one message outperformed another.",
                    ],
                )
            rows = [
                ChartRow(
                    label=variant.label or variant.variant_key,
                    mean=variant.objective_rate.mean,
                    lower=variant.objective_rate.lower,
                    upper=variant.objective_rate.upper,
                    n=variant.objective_rate.n,
                )
                for variant in scoreboard.variants
                if variant.objective_rate.n > 0
            ]
            if rows:
                _chart_slide(
                    prs, "Objective rate by message arena",
                    render_interval_rows_png(
                        rows, "Objective rate by message arena",
                        domain=(0.0, 1.0), signed=False,
                        axis_label="Share of active agents taking the objective action",
                    ),
                )

        arc = [
            ChartRow(
                label=str(point.round_number),
                mean=point.valence.mean,
                lower=point.valence.lower,
                upper=point.valence.upper,
                n=point.valence.n,
            )
            for point in analysis.sentiment_timeline
            if point.valence.n > 0
        ]
        if len(arc) >= 2:
            _chart_slide(
                prs, "Measured valence by round", render_sentiment_arc_png(arc)
            )

        platforms = _interval_rows(
            analysis.by_platform, lambda s: platform_label(s.platform)
        )
        if len(platforms) >= 2:
            _chart_slide(
                prs, "Measured valence by platform",
                render_interval_rows_png(platforms, "Measured valence by platform"),
            )

        archetypes = _interval_rows(analysis.by_archetype, lambda s: s.archetype)
        if len(archetypes) >= 2:
            _chart_slide(
                prs, "Measured valence by archetype",
                render_interval_rows_png(archetypes, "Measured valence by archetype"),
            )

        cohorts = _interval_rows(
            analysis.by_cohort, lambda s: COHORT_NAMES.get(s.cohort, s.cohort)
        )
        if len(cohorts) >= 2:
            _chart_slide(
                prs, "Buyers against the incumbent-aligned cohort",
                render_interval_rows_png(
                    cohorts, "Buyers against the incumbent-aligned cohort"
                ),
            )

        if analysis.objections:
            _bullet_slide(
                prs, "Objections, by load-bearing weight",
                [
                    f"{objection.label} — weight {objection.load_bearing_score:.1f}, "
                    f"{objection.agent_count} agents"
                    + (
                        f", first seen R{objection.first_round_seen}"
                        if objection.first_round_seen is not None
                        else ""
                    )
                    for objection in analysis.objections[:6]
                ],
                note="Ranked by reach × intensity × cohort spread, not by how "
                "often an objection was repeated.",
            )

        # PRD §4 — its own slide, before methodology. A deck is presented one
        # slide at a time, and a disclosure sharing a slide with the platform
        # list is a disclosure that gets skipped past.
        if analysis.adversarial.enabled:
            bullets = [analysis.adversarial.disclosure]
            if analysis.adversarial.roles:
                bullets.append(
                    "Roles: "
                    + ", ".join(
                        f"{key.replace('_', ' ')} ({value})"
                        for key, value in sorted(analysis.adversarial.roles.items())
                    )
                )
            _bullet_slide(prs, "Adversarial cohort disclosure", bullets)

    for section in sections:
        content = clean_report_output(section.get("content") or "")
        bullets = [line.strip() for line in content.split("\n") if line.strip()][:5]
        _bullet_slide(prs, section.get("title") or "Findings", bullets)

    _bullet_slide(
        prs, "Method",
        [
            f"Question: {simulation.get('prediction_goal') or '—'}",
            "Platforms: "
            + (", ".join(platform_label(p) for p in (simulation.get("platforms") or [])) or "—"),
            f"Rounds: {simulation.get('max_rounds') or '—'}",
            f"Message arenas: {simulation.get('variants') or 1}",
            "Every agent in this run is synthetic. Confidence intervals are "
            "computed across agents; unmeasured values are absent, not zero.",
        ],
    )

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    result = buf.read()
    logger.info(
        "pptx_exported",
        report_id=report_id,
        bytes=len(result),
        slides=len(prs.slides),
        measured=analysis is not None,
    )
    return result
